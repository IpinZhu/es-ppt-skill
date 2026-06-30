"""Linux-friendly PPTX generator for the es-ppt-skill template (no PowerPoint required).

Phase 2 engine. Produces a fully editable .pptx by:
  * cloning content template slides once per JSON entry
  * reordering so layout becomes [cover, toc, ...content..., thanks]
  * filling every text frame's first paragraph in place (preserving pPr / first-run rPr)

Supports 7 content layouts:
  gallery, highlight, dual, columns, text-image, table, table-image

Math support
------------
Every text we inject is scanned for inline LaTeX delimited by ``$ ... $``.
Each math segment is converted with latex2mathml -> mathml2omml and embedded
into the paragraph as ``<a14:m><m:oMath>...</m:oMath></a14:m>`` so PowerPoint /
WPS / Keynote render it as a native equation, not as raw ``$...$`` text.
Plain text segments keep using ``<a:r>`` runs and inherit the original rPr.
"""
from __future__ import annotations
import sys
import os
import json
import shutil
import copy
import re
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn, nsmap
from pptx.opc.package import PackURI
from pptx.parts.slide import SlidePart
from lxml import etree


PRES_SLIDE_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"

# Namespaces used for embedding OMML inside DrawingML paragraphs.
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"


def _qn(prefix: str, tag: str, ns_uri: str) -> str:
    return f"{{{ns_uri}}}{tag}"


# ---------- LaTeX -> OMML ----------
try:
    import latex2mathml.converter as _l2m
    import mathml2omml as _m2o
    HAS_MATH = True
except Exception:  # pragma: no cover
    HAS_MATH = False


_MATH_RE = re.compile(r"\$([^$\n]+)\$")


def latex_to_omml_element(latex_src: str):
    """Return an lxml Element for ``<m:oMath>`` representing *latex_src*.

    Returns None if conversion fails or math libs are missing — callers must
    fall back to plain text in that case.
    """
    if not HAS_MATH:
        return None
    try:
        mml = _l2m.convert(latex_src)
        omml = _m2o.convert(mml)
    except Exception:
        return None
    wrapped = (
        f'<root xmlns:m="{M_NS}">{omml}</root>'
    )
    try:
        root = etree.fromstring(wrapped.encode("utf-8"))
    except etree.XMLSyntaxError:
        return None
    omath_nodes = root.findall(_qn("m", "oMath", M_NS))
    if not omath_nodes:
        return None
    return omath_nodes[0]


def split_text_with_math(text: str):
    """Yield ``("text", str)`` and ``("math", latex_src)`` segments."""
    if not text:
        return
    pos = 0
    for m in _MATH_RE.finditer(text):
        if m.start() > pos:
            yield "text", text[pos:m.start()]
        yield "math", m.group(1)
        pos = m.end()
    if pos < len(text):
        yield "text", text[pos:]


# ---------- slide cloning / reordering ----------
def clone_slide(prs, source_index):
    src_part = prs.slides[source_index].part
    pkg = prs.part.package
    used = {p.partname.lower() for p in pkg.iter_parts()}
    i = 1
    while True:
        candidate = PackURI(f"/ppt/slides/slide{i}.xml")
        if str(candidate).lower() not in used:
            break
        i += 1
    new_element = copy.deepcopy(src_part._element)
    new_part = SlidePart(candidate, src_part.content_type, pkg, new_element)
    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_part.rels.get_or_add(rel.reltype, rel.target_part)
    rId = prs.part.relate_to(new_part, PRES_SLIDE_RELTYPE)
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(e.get("id")) for e in sldIdLst.findall(qn("p:sldId"))]
    new_id = max(existing_ids + [255]) + 1
    sldId = etree.SubElement(sldIdLst, qn("p:sldId"))
    sldId.set("id", str(new_id))
    sldId.set(qn("r:id"), rId)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)
    return prs.slides[len(prs.slides) - 1]


def remove_slide(prs, slide_index):
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst.findall(qn("p:sldId")))
    target = entries[slide_index]
    rId = target.get(qn("r:id"))
    sldIdLst.remove(target)
    prs.part.drop_rel(rId)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)


def move_slide(prs, src_index, dst_index):
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst.findall(qn("p:sldId")))
    target = entries[src_index]
    sldIdLst.remove(target)
    entries = list(sldIdLst.findall(qn("p:sldId")))
    if dst_index >= len(entries):
        sldIdLst.append(target)
    else:
        sldIdLst.insert(dst_index, target)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)


# ---------- text helpers ----------
def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def _capture_first_run_rPr(text_frame):
    paras = list(text_frame.paragraphs)
    if not paras:
        return None, None
    first = paras[0]
    pPr = first._p.find(qn("a:pPr"))
    pPr_copy = copy.deepcopy(pPr) if pPr is not None else None
    rPr_copy = None
    for r in first.runs:
        rPr_node = r._r.find(qn("a:rPr"))
        if rPr_node is not None:
            rPr_copy = copy.deepcopy(rPr_node)
            break
    return pPr_copy, rPr_copy


_ILLEGAL_XML_CHARS_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1F\x7F-\x9F]')


def _sanitize_xml(text):
    return _ILLEGAL_XML_CHARS_RE.sub('', text)


def _add_run_with_text(p_elem, text, rPr_copy):
    r = etree.SubElement(p_elem, qn("a:r"))
    if rPr_copy is not None:
        r.append(copy.deepcopy(rPr_copy))
    t = etree.SubElement(r, qn("a:t"))
    t.text = _sanitize_xml(text)


def _add_math(p_elem, latex_src):
    """Append <a14:m><m:oMath/></a14:m> to *p_elem*. Returns True on success."""
    omath = latex_to_omml_element(latex_src)
    if omath is None:
        return False
    a14_m = etree.SubElement(p_elem, _qn("a14", "m", A14_NS), nsmap={"a14": A14_NS, "m": M_NS})
    a14_m.append(omath)
    return True


def _fill_paragraph(p_elem, text, rPr_copy):
    """Populate paragraph *p_elem* with mixed text + math runs based on *text*.

    A leading <a:pPr> child, if any, is preserved.
    Existing <a:r>, <a:br>, <a:fld>, <a14:m> children are removed first.
    """
    for child in list(p_elem):
        if child.tag in (qn("a:r"), qn("a:br"), qn("a:fld"), _qn("a14", "m", A14_NS)):
            p_elem.remove(child)
    for kind, payload in split_text_with_math(text):
        if kind == "text":
            if payload:
                _add_run_with_text(p_elem, payload, rPr_copy)
        else:  # math
            if not _add_math(p_elem, payload):
                _add_run_with_text(p_elem, f"${payload}$", rPr_copy)


def set_text_keep_format(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    p = etree.SubElement(txBody, qn("a:p"))
    if pPr_copy is not None:
        p.append(copy.deepcopy(pPr_copy))
    _fill_paragraph(p, new_text, rPr_copy)


def set_multiline_text(shape, lines, line_height_emu=457200, min_lines=3):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    for line in lines:
        p = etree.SubElement(txBody, qn("a:p"))
        if pPr_copy is not None:
            p.append(copy.deepcopy(pPr_copy))
        _fill_paragraph(p, line, rPr_copy)
    # auto-fit shape height to number of lines (only grow, never shrink)
    n = max(len(lines), min_lines)
    desired = n * line_height_emu
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            ext = xfrm.find(qn("a:ext"))
            if ext is not None:
                try:
                    current_cy = int(ext.get("cy", 0))
                    if desired > current_cy:
                        ext.set("cy", str(desired))
                except (ValueError, TypeError):
                    pass
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for tag in ("a:noAutofit", "a:normAutofit", "a:spAutoFit"):
            for n_ in bodyPr.findall(qn(tag)):
                bodyPr.remove(n_)
        etree.SubElement(bodyPr, qn("a:noAutofit"))
        bodyPr.set("wrap", "square")


def set_two_line_with_break(shape, line1, line2):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    p = etree.SubElement(txBody, qn("a:p"))
    if pPr_copy is not None:
        p.append(copy.deepcopy(pPr_copy))
    _fill_paragraph(p, line1, rPr_copy)
    etree.SubElement(p, qn("a:br"))
    tmp = etree.SubElement(txBody, qn("a:p"))
    _fill_paragraph(tmp, line2, rPr_copy)
    for child in list(tmp):
        if child.tag == qn("a:pPr"):
            continue
        p.append(child)
    txBody.remove(tmp)


# ---------- image helpers ----------
def _collect_image_placeholders(slide):
    """Return list of placeholder shapes sorted by left position."""
    placeholders = []
    for shape in slide.shapes:
        if shape.shape_type == 1 and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith("图片"):
                placeholders.append(shape)
    placeholders.sort(key=lambda s: s.left)
    return placeholders


def replace_images_in_slide(slide, images):
    """Replace image placeholders in slide with actual images.
    images: list of file paths. Placeholders matched left-to-right.
    """
    if not images:
        return
    placeholders = _collect_image_placeholders(slide)
    for ph, img_path in zip(placeholders, images):
        if not img_path or not os.path.exists(img_path):
            continue
        left, top, width, height = ph.left, ph.top, ph.width, ph.height
        sp = ph._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(img_path, left, top, width, height)


def remove_image_placeholders(slide):
    """Remove all image placeholder shapes from slide."""
    for shape in list(slide.shapes):
        if shape.shape_type == 1 and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith("图片"):
                sp = shape._element
                sp.getparent().remove(sp)


# ---------- table helpers ----------
def _find_table_shape(slide):
    for shape in iter_shapes(slide.shapes):
        if shape.shape_type == 19 and hasattr(shape, "table"):
            return shape
    return None


def _capture_cell_rPr(table):
    """Capture rPr from the first cell's first run as template."""
    if len(table.rows) == 0 or len(table.columns) == 0:
        return None
    cell = table.cell(0, 0)
    for r in cell.text_frame.paragraphs[0].runs:
        rPr_node = r._r.find(qn("a:rPr"))
        if rPr_node is not None:
            return copy.deepcopy(rPr_node)
    return None


def fill_table_shape(table_shape, table_data):
    """Fill a table shape with data. table_data = {"headers": [...], "rows": [[...], ...]}"""
    if table_shape is None or not hasattr(table_shape, "table"):
        return
    table = table_shape.table
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers and not rows:
        return

    needed_rows = 1 + len(rows)
    needed_cols = max(len(headers), max((len(r) for r in rows), default=0))

    # Adjust rows via low-level XML
    def _add_row():
        tbl = table._tbl
        trs = tbl.findall(qn("a:tr"))
        if not trs:
            return
        last_tr = trs[-1]
        new_tr = copy.deepcopy(last_tr)
        for tc in new_tr.findall(qn("a:tc")):
            txBody = tc.find(qn("a:txBody"))
            if txBody is not None:
                for p in txBody.findall(qn("a:p")):
                    for child in list(p):
                        if child.tag != qn("a:pPr"):
                            p.remove(child)
                    if len(p) == 0:
                        r = etree.SubElement(p, qn("a:r"))
                        t = etree.SubElement(r, qn("a:t"))
                        t.text = ""
        tbl.append(new_tr)

    def _add_col():
        tbl = table._tbl
        for tr in tbl.findall(qn("a:tr")):
            tcs = tr.findall(qn("a:tc"))
            if not tcs:
                continue
            last_tc = tcs[-1]
            new_tc = copy.deepcopy(last_tc)
            txBody = new_tc.find(qn("a:txBody"))
            if txBody is not None:
                for p in txBody.findall(qn("a:p")):
                    for child in list(p):
                        if child.tag != qn("a:pPr"):
                            p.remove(child)
                    if len(p) == 0:
                        r = etree.SubElement(p, qn("a:r"))
                        t = etree.SubElement(r, qn("a:t"))
                        t.text = ""
            tr.append(new_tc)
        tblGrid = tbl.find(qn("a:tblGrid"))
        if tblGrid is not None:
            last_gridCol = tblGrid.findall(qn("a:gridCol"))[-1]
            new_gridCol = copy.deepcopy(last_gridCol)
            tblGrid.append(new_gridCol)

    while len(table.rows) < needed_rows:
        _add_row()

    while len(table.columns) < needed_cols:
        _add_col()

    # Capture rPr from first cell for math rendering
    rPr_copy = _capture_cell_rPr(table)

    # Fill header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        p = cell.text_frame.paragraphs[0]._p
        _fill_paragraph(p, str(h), rPr_copy)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True

    # Fill data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j < needed_cols:
                cell = table.cell(i + 1, j)
                p = cell.text_frame.paragraphs[0]._p
                _fill_paragraph(p, str(val), rPr_copy)

    # Clear extra cells in header / data range
    for i in range(needed_rows):
        for j in range(needed_cols, len(table.columns)):
            table.cell(i, j).text = ""


# ---------- layout detection helpers ----------
def _is_title_text(txt):
    """Heuristic to detect title placeholder like '01 模板一'."""
    if not txt:
        return False
    return txt[0].isdigit() and "模板" in txt


def _fill_title_and_subtitle(slide, s_data):
    """Fill title and subtitle on any content slide.
    Returns a set of shape _elements that were handled."""
    title = s_data.get("title", "")
    subtitle = s_data.get("subtitle", "")
    handled = set()
    arrow_shapes = []
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if _is_title_text(txt):
            set_text_keep_format(shape, title)
            handled.add(shape._element)
        elif txt.startswith("➢"):
            # Defer arrow shapes; only the topmost one is the page subtitle.
            # Others (e.g. column subtitles like "➢ 小标题2") are left to layout fillers.
            arrow_shapes.append((shape, shape.top))
    if arrow_shapes:
        arrow_shapes.sort(key=lambda x: x[1])
        set_text_keep_format(arrow_shapes[0][0], subtitle)
        handled.add(arrow_shapes[0][0]._element)
    return handled


# ---------- layout fillers ----------
def fill_gallery(slide, s_data):
    handled = _fill_title_and_subtitle(slide, s_data)
    for shape in iter_shapes(slide.shapes):
        if shape._element in handled or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "列表1" in txt:
            bullets = s_data.get("bullets", [])
            set_multiline_text(shape, bullets)
    replace_images_in_slide(slide, s_data.get("images", []))


def fill_highlight(slide, s_data):
    handled = _fill_title_and_subtitle(slide, s_data)
    for shape in iter_shapes(slide.shapes):
        if shape._element in handled or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "*高亮" in txt or "HIGHLIGHT" in txt.upper():
            set_text_keep_format(shape, s_data.get("highlight", ""))
        elif "列表1" in txt:
            bullets = s_data.get("bullets", [])
            set_multiline_text(shape, bullets)
    replace_images_in_slide(slide, s_data.get("images", []))


def fill_dual(slide, s_data):
    handled = _fill_title_and_subtitle(slide, s_data)
    blocks = s_data.get("blocks", [])
    summary = s_data.get("summary", "")
    # Collect candidate text shapes (exclude already-handled title/subtitle)
    candidates = []
    for shape in iter_shapes(slide.shapes):
        if shape._element in handled or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        candidates.append((shape.top, shape, txt))
    candidates.sort(key=lambda x: x[0])
    # Fill blocks into first N candidates, summary into last
    for idx, (_, shape, txt) in enumerate(candidates):
        if idx < len(blocks):
            lines = blocks[idx].split("\n")
            set_multiline_text(shape, lines, min_lines=1)
        elif idx == len(candidates) - 1 and summary:
            set_text_keep_format(shape, summary)
        else:
            set_text_keep_format(shape, "")


def fill_columns(slide, s_data):
    handled = _fill_title_and_subtitle(slide, s_data)
    columns = s_data.get("columns", [])
    # Collect remaining text shapes
    subtitle_infos = []
    content_infos = []
    for shape in iter_shapes(slide.shapes):
        if shape._element in handled or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt.startswith("➢"):
            subtitle_infos.append({"shape": shape, "top": shape.top})
        elif "介绍" in txt or txt.startswith("方法") or "内容" in txt:
            content_infos.append({"shape": shape, "top": shape.top})

    subtitle_infos.sort(key=lambda x: x["top"])
    content_infos.sort(key=lambda x: x["top"])

    for idx, info in enumerate(subtitle_infos):
        if idx < len(columns):
            set_text_keep_format(info["shape"], columns[idx].get("subtitle", ""))
        else:
            set_text_keep_format(info["shape"], "")

    for idx, info in enumerate(content_infos):
        if idx < len(columns):
            content = columns[idx].get("content", "")
            lines = content.split("\n")
            set_multiline_text(info["shape"], lines, min_lines=1)
        else:
            set_text_keep_format(info["shape"], "")

    replace_images_in_slide(slide, s_data.get("images", []))


def fill_text_image(slide, s_data):
    handled = _fill_title_and_subtitle(slide, s_data)
    text = s_data.get("text", "")
    for shape in iter_shapes(slide.shapes):
        if shape._element in handled or not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if "描述性文字" in txt:
            lines = text.split("\n")
            set_multiline_text(shape, lines, min_lines=1)
    replace_images_in_slide(slide, [s_data.get("image", "")] if s_data.get("image") else [])


def fill_table_image(slide, s_data):
    _fill_title_and_subtitle(slide, s_data)
    table_data = s_data.get("table")
    if table_data:
        table_shape = _find_table_shape(slide)
        fill_table_shape(table_shape, table_data)
    replace_images_in_slide(slide, [s_data.get("image", "")] if s_data.get("image") else [])


def fill_table(slide, s_data):
    _fill_title_and_subtitle(slide, s_data)
    table_data = s_data.get("table")
    if table_data:
        table_shape = _find_table_shape(slide)
        fill_table_shape(table_shape, table_data)
    # Remove image placeholder for pure-table layout
    remove_image_placeholders(slide)


LAYOUT_FILLERS = {
    "gallery": fill_gallery,
    "highlight": fill_highlight,
    "dual": fill_dual,
    "columns": fill_columns,
    "text-image": fill_text_image,
    "table": fill_table,
    "table-image": fill_table_image,
}

LAYOUT_MAP = {
    "gallery": 2,
    "highlight": 3,
    "dual": 4,
    "columns": 5,
    "text-image": 6,
    "table": 8,
    "table-image": 8,
}


# ---------- standard fillers ----------
def fill_cover(slide, cover):
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "大大大大标题" in text or "子标题" in text:
            set_two_line_with_break(shape, cover.get("title", ""), cover.get("subtitle", ""))
        elif ("汇报人" in text or "申 请 人" in text) and ("导" in text and "师" in text):
            lines = [
                f"汇 报 人 ：{cover.get('reporter','')}",
                f"指导老师：{cover.get('instructor','')}",
                f"专 业 ：{cover.get('major','')}",
                f"日 期 ：{cover.get('date','')}",
            ]
            set_multiline_text(shape, lines)
        elif text.strip().endswith("年5月") or text.strip().startswith("2025年") or text.strip().startswith("2026年"):
            set_text_keep_format(shape, cover.get("date", ""))


def fill_toc(slide, toc):
    shapes = list(iter_shapes(slide.shapes))
    num_shapes = []
    title_shapes = []
    for sh in shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        # Skip decorative / non-entry texts
        if t in ("目 录", "CONTENTS", "大工至善 大学至真"):
            continue
        if t.isdigit() and len(t) <= 2:
            continue
        # Match two-digit entry numbers like "01", "02"
        if re.match(r"^\d{2}$", t):
            num_shapes.append(sh)
        elif t:
            title_shapes.append(sh)
    num_shapes.sort(key=lambda s: s.top)
    title_shapes.sort(key=lambda s: s.top)
    n = min(len(toc), len(num_shapes), len(title_shapes))
    for i in range(n):
        set_text_keep_format(num_shapes[i], f"{i+1:02d}")
        set_text_keep_format(title_shapes[i], toc[i])


# ---------- post-processing ----------
def _hoist_math_namespaces(pptx_path):
    """Move xmlns:a14 / xmlns:m declarations from each <a14:m> up to the root
    <p:sld>, so the file is smaller and matches the layout PowerPoint emits."""
    import zipfile
    import re as _re
    A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
    M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        items = {n: zin.read(n) for n in zin.namelist()}
    changed = False
    for name in list(items):
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        raw = items[name].decode("utf-8")
        if "<a14:m " not in raw:
            continue
        def add_ns(m):
            head = m.group(1)
            new = head
            if "xmlns:a14=" not in new:
                new += f' xmlns:a14="{A14}"'
            if "xmlns:m=" not in new:
                new += f' xmlns:m="{M}"'
            return new + ">"
        raw = _re.sub(r"(<p:sld\b[^>]*?)>", add_ns, raw, count=1)
        raw = raw.replace(f'<a14:m xmlns:a14="{A14}" xmlns:m="{M}">', "<a14:m>")
        items[name] = raw.encode("utf-8")
        changed = True
    if not changed:
        return
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, b in items.items():
            zout.writestr(n, b)


# ---------- main ----------
def process(json_path, template_path, output_path):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    shutil.copyfile(template_path, output_path)
    prs = Presentation(output_path)
    slides_data = data.get("slides", [])

    # Clone required template slides
    for s in slides_data:
        layout = s.get("layout", "gallery")
        idx = LAYOUT_MAP.get(layout, 2)
        clone_slide(prs, idx)
        last = len(prs.slides) - 1
        move_slide(prs, last, last - 1)

    # Remove all original content template slides (indices 2-8) from high to low
    for idx in sorted([8, 7, 6, 5, 4, 3, 2], reverse=True):
        if idx < len(prs.slides):
            remove_slide(prs, idx)

    fill_cover(prs.slides[0], data.get("cover", {}))
    fill_toc(prs.slides[1], data.get("toc", []))
    for i, s_data in enumerate(slides_data):
        slide = prs.slides[2 + i]
        layout = s_data.get("layout", "gallery")
        filler = LAYOUT_FILLERS.get(layout, fill_gallery)
        filler(slide, s_data)

    prs.save(output_path)
    _hoist_math_namespaces(output_path)
    print(f"Successfully generated {output_path}")
    if not HAS_MATH:
        print("Note: latex2mathml / mathml2omml not installed — math fell back to plain text. "
              "Run: pip install latex2mathml mathml2omml")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python build_pptx_linux.py <data.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    process(sys.argv[1], sys.argv[2], sys.argv[3])
