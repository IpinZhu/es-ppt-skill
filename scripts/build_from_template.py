import sys
import json
import os
import shutil
import win32com.client
from pptx import Presentation

# Import shared fillers and helpers from the Linux engine
# (they are pure python-pptx operations, cross-platform)
from build_pptx_linux import (
    iter_shapes,
    set_text_keep_format,
    set_multiline_text,
    set_two_line_with_break,
    replace_images_in_slide,
    remove_image_placeholders,
    fill_table_shape,
    _find_table_shape,
    fill_gallery,
    fill_highlight,
    fill_dual,
    fill_columns,
    fill_text_image,
    fill_table,
    fill_table_image,
    fill_cover,
    fill_toc,
    LAYOUT_FILLERS,
)

# 1-based slide indices for the Windows COM engine
WIN32_LAYOUT_MAP = {
    "gallery":      3,   # Slide 2 (0-based)
    "highlight":    4,   # Slide 3
    "dual":         5,   # Slide 4
    "columns":      6,   # Slide 5
    "text-image":   7,   # Slide 6
    "table":        8,   # Slide 7 (table-image template, image removed)
    "table-image":  8,   # Slide 7
}


def prepare_slides_win32(pptx_path, slides_data):
    abs_path = os.path.abspath(pptx_path)
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(abs_path, WithWindow=False)
    try:
        for s_data in slides_data:
            layout = s_data.get('layout', 'gallery')
            idx = WIN32_LAYOUT_MAP.get(layout, 3)
            new_slide_range = presentation.Slides(idx).Duplicate()
            new_slide = new_slide_range(1)
            # Move the new slide to just before the thanks slide
            new_slide.MoveTo(presentation.Slides.Count - 1)

        # Delete original content template slides (indices 3~8, 1-based) from high to low.
        # Clones have been moved to the end, so indices 3~8 are the originals.
        for idx in range(8, 2, -1):
            presentation.Slides(idx).Delete()

        presentation.Save()
    except Exception as e:
        print(f"Error duplicating slides: {e}")
    finally:
        presentation.Close()
        if powerpoint.Presentations.Count == 0:
            powerpoint.Quit()


def process_presentation(json_path, template_path, output_path):
    shutil.copyfile(template_path, output_path)

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    slides_data = data.get('slides', [])
    if len(slides_data) == 0:
        print("No content slides to generate.")
        return

    # Phase 1: Duplicate slides using win32com
    prepare_slides_win32(output_path, slides_data)

    # Phase 2: Populate text using python-pptx
    prs = Presentation(output_path)

    fill_cover(prs.slides[0], data.get('cover', {}))
    fill_toc(prs.slides[1], data.get('toc', []))

    for i, s_data in enumerate(slides_data):
        slide = prs.slides[i + 2]
        layout = s_data.get('layout', 'gallery')
        filler = LAYOUT_FILLERS.get(layout, fill_gallery)
        filler(slide, s_data)

    prs.save(output_path)
    print(f"Successfully generated {output_path}")


if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: python build_from_template.py <data.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    process_presentation(sys.argv[1], sys.argv[2], sys.argv[3])
